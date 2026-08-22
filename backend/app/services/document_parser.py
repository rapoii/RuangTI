import os
import zipfile
import csv
import gzip
import re
import logging
from typing import Dict

logger = logging.getLogger("RuangTI.DocParser")

MAX_TEXT_CHAR_LIMIT = 40000  # Safe token budget for attached documents

# ---------------------------------------------------------------------------
# 1. Plain Text, Markdown & Code Files
# ---------------------------------------------------------------------------
def parse_text_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Reads a plain text, code, json, yaml, sql, or markdown file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read(max_chars)
        return content
    except Exception as e:
        logger.error(f"Error reading text file {file_path}: {e}")
        return f"[Gagal membaca file teks: {e}]"

# ---------------------------------------------------------------------------
# 2. Office Documents (Word, Excel, PDF, CSV)
# ---------------------------------------------------------------------------
def parse_docx_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts text, headings, and tables from Word (.docx) documents."""
    try:
        import docx
        doc = docx.Document(file_path)
        parts = []
        
        # Paragraphs & Headings
        for p in doc.paragraphs:
            text = p.text.strip()
            if text:
                if p.style and p.style.name.startswith("Heading"):
                    parts.append(f"\n### {text}\n")
                else:
                    parts.append(text)
                    
        # Tables
        for t_idx, table in enumerate(doc.tables, 1):
            parts.append(f"\n[Tabel {t_idx}]")
            for row in table.rows:
                row_vals = [c.text.strip().replace("\n", " ") for c in row.cells]
                parts.append("| " + " | ".join(row_vals) + " |")
                
        full_text = "\n".join(parts)
        return full_text[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing docx {file_path}: {e}")
        return f"[Gagal mengekstrak dokumen Word: {e}]"

def parse_excel_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts sheet names, headers, and row samples from Excel (.xlsx/.xls) workbooks."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"\n=== Sheet: {sheet_name} ===")
            
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count > 100:
                    parts.append("... [Baris data lainnya dipersingkat]")
                    break
                if any(v is not None for v in row):
                    row_vals = [str(v).strip() if v is not None else "" for v in row]
                    parts.append("| " + " | ".join(row_vals[:25]) + " |")
                    row_count += 1
                    
        full_text = "\n".join(parts)
        return full_text[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing excel {file_path}: {e}")
        return f"[Gagal mengekstrak data spreadsheet Excel: {e}]"

def parse_csv_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Parses CSV text data with table formatting."""
    try:
        parts = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for idx, row in enumerate(reader):
                if idx > 150:
                    parts.append("... [Baris CSV selanjutnya dipersingkat]")
                    break
                parts.append("| " + " | ".join(row[:25]) + " |")
        full_text = "\n".join(parts)
        return full_text[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing CSV {file_path}: {e}")
        return f"[Gagal membaca file CSV: {e}]"

def parse_pdf_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts plain text and tabular data from PDF documents using PyMuPDF (fitz) or pypdf fallback."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        parts = []
        for page_idx, page in enumerate(doc):
            text = page.get_text("text").strip()
            if text:
                parts.append(f"\n--- Halaman {page_idx + 1} ---\n{text}")
            if sum(len(p) for p in parts) >= max_chars:
                break
        full_text = "\n".join(parts)
        return full_text[:max_chars]
    except Exception as e_fitz:
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            parts = []
            for page_idx, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    parts.append(f"\n--- Halaman {page_idx + 1} ---\n{text.strip()}")
                if sum(len(p) for p in parts) >= max_chars:
                    break
            full_text = "\n".join(parts)
            return full_text[:max_chars]
        except Exception as e_pypdf:
            logger.error(f"Error parsing PDF {file_path}: fitz({e_fitz}), pypdf({e_pypdf})")
            return f"[Gagal membaca dokumen PDF: {e_fitz}]"

def parse_pptx_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts text, slide titles, bullet points, and tables from PowerPoint (.pptx) presentations."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = ["=== ANALISIS PRESENTASI POWERPOINT (.PPTX) ==="]
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_parts = [f"\n--- Slide {slide_idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(f"- {text}")
                elif shape.has_table:
                    slide_parts.append("[Tabel Slide]")
                    for row in shape.table.rows:
                        row_vals = [c.text.strip().replace("\n", " ") for c in row.cells]
                        slide_parts.append("| " + " | ".join(row_vals) + " |")
                        
            parts.append("\n".join(slide_parts))
            if sum(len(p) for p in parts) >= max_chars:
                break
                
        return "\n".join(parts)[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing PPTX {file_path}: {e}")
        return f"[Gagal membaca file presentasi PowerPoint: {e}]"

# ---------------------------------------------------------------------------
# 3. Archive Files (ZIP, TAR, GZ, 7Z)
# ---------------------------------------------------------------------------
def parse_zip_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts file list tree and contents of key code/text files from a ZIP archive."""
    try:
        with zipfile.ZipFile(file_path, "r") as z:
            namelist = z.namelist()
            parts = [f"=== Struktur Berkas ZIP ({len(namelist)} file) ==="]
            
            for name in namelist[:60]:
                parts.append(f"- {name}")
            if len(namelist) > 60:
                parts.append(f"... dan {len(namelist) - 60} file lainnya.")
                
            parts.append("\n=== Cuplikan Berkas Penting di Dalam ZIP ===")
            read_budget = max_chars - sum(len(p) for p in parts)
            
            for name in namelist:
                if read_budget <= 1000:
                    break
                ext = name.split(".")[-1].lower() if "." in name else ""
                if ext in ["py", "js", "ts", "tsx", "jsx", "json", "yaml", "yml", "md", "txt", "sql", "csv", "html", "css", "dxf", "step", "stp", "gcode"]:
                    try:
                        with z.open(name) as zf:
                            content = zf.read(8000).decode("utf-8", errors="replace")
                            parts.append(f"\nFile: `{name}`:\n```{ext}\n{content}\n```")
                            read_budget -= len(content)
                    except Exception:
                        pass
                        
            full_text = "\n".join(parts)
            return full_text[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing ZIP {file_path}: {e}")
        return f"[Gagal mengekstrak berkas arsip ZIP: {e}]"

# ---------------------------------------------------------------------------
# 4. AutoCAD 2D CAD Files (.dwg, .dxf)
# ---------------------------------------------------------------------------
def parse_dwg_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts layers, modelspace entities, dimensions, and text annotations from AutoCAD .dwg files."""
    try:
        import ezdwg
        doc = ezdwg.read(file_path)
        parts = ["=== ANALISIS GAMBAR TEKNIK AUTOCAD (DWG) ==="]
        
        # Units and metadata
        units = getattr(doc, "units", "Unknown")
        parts.append(f"- Satuan Ukuran (Units): {units if units else 'Milimeter / Unit Gambar'}")
        
        # Header extents
        header = doc.header_variables() if hasattr(doc, "header_variables") else {}
        if "extmin" in header and "extmax" in header:
            parts.append(f"- Batas Area Gambar (Extents): {header['extmin']} s.d. {header['extmax']}")
            
        # Modelspace entity query
        msp = doc.modelspace()
        entity_counts: Dict[str, int] = {}
        text_annotations = []
        layers = set()
        
        for e in msp.query("LINE LWPOLYLINE ARC CIRCLE ELLIPSE POINT TEXT MTEXT DIMENSION INSERT MINSERT HATCH SPLINE"):
            dxftype = getattr(e, "dxftype", "UNKNOWN")
            entity_counts[dxftype] = entity_counts.get(dxftype, 0) + 1
            
            dxf = getattr(e, "dxf", {})
            if isinstance(dxf, dict):
                layer = dxf.get("layer")
                if layer:
                    layers.add(str(layer))
                # Extract text
                if dxftype in ["TEXT", "MTEXT"]:
                    txt = dxf.get("text") or dxf.get("plain_text")
                    if txt and str(txt).strip():
                        text_annotations.append(str(txt).strip())
                elif dxftype == "DIMENSION":
                    dim_txt = dxf.get("text")
                    if dim_txt and str(dim_txt).strip():
                        text_annotations.append(f"[Dimensi: {dim_txt}]")
                        
        # Summary counts
        parts.append("\n[RINGKASAN ELEMEN GAMBAR]:")
        for k, v in sorted(entity_counts.items(), key=lambda x: -x[1]):
            parts.append(f"- {k}: {v} entitas")
            
        if layers:
            parts.append(f"\n[DAFTAR LAYER ({len(layers)})]:")
            parts.append(", ".join(sorted(list(layers))[:20]))
            
        if text_annotations:
            parts.append(f"\n[TEKS & ETIKET GAMBAR ({len(text_annotations)})]:")
            for t in text_annotations[:25]:
                parts.append(f"- {t}")
                
        return "\n".join(parts)[:max_chars]
    except Exception as e:
        logger.warning(f"ezdwg failed for {file_path}, falling back to string extract: {e}")
        # Fallback binary string extraction
        try:
            with open(file_path, "rb") as f:
                raw = f.read(1000000)
            strings = re.findall(rb'[A-Za-z0-9_ \-\.\(\)\,\:\/\\><\=\+\*\[\]]{4,}', raw)
            decoded = [s.decode('latin1', errors='ignore').strip() for s in strings if len(s.strip()) >= 4]
            texts = [s for s in decoded if any(k in s.lower() for k in ['untirta', 'judul', 'skala', 'revisi', 'layer', 'ptlf', 'layout', 'mesin', 'dimensi', 'part'])]
            return f"=== ANALISIS GAMBAR AUTOCAD DWG (FALLBACK SCAN) ===\n- File: {os.path.basename(file_path)}\n- Deteksi Teks Etiket: {', '.join(texts[:20]) if texts else 'Format DWG biner terenkripsi'}"
        except Exception:
            return f"[Gagal membaca file AutoCAD DWG: {e}]"

def parse_dxf_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts layers, entities, dimensions, and text annotations from Drawing Interchange Format (.dxf) files."""
    try:
        import ezdxf
        doc = ezdxf.readfile(file_path)
        parts = ["=== ANALISIS GAMBAR CAD DXF (Drawing Interchange Format) ==="]
        
        # Header Info
        header = doc.header
        if "$INSUNITS" in header:
            unit_map = {0: "Unitless", 1: "Inches", 2: "Feet", 4: "Millimeters", 5: "Centimeters", 6: "Meters"}
            unit_code = header["$INSUNITS"]
            parts.append(f"- Satuan Gambar: {unit_map.get(unit_code, f'Code {unit_code}')}")
            
        # Layers
        layers = [layer.dxf.name for layer in doc.layers]
        parts.append(f"- Total Layer: {len(layers)} ({', '.join(layers[:15])})")
        
        # Modelspace entities
        msp = doc.modelspace()
        entity_counts: Dict[str, int] = {}
        texts = []
        
        for e in msp:
            t = e.dxftype()
            entity_counts[t] = entity_counts.get(t, 0) + 1
            if t in ["TEXT", "MTEXT"]:
                val = e.dxf.text if hasattr(e.dxf, "text") else ""
                if val:
                    texts.append(val.strip())
            elif t == "DIMENSION":
                val = getattr(e.dxf, "text", "")
                if val:
                    texts.append(f"[Dimensi: {val}]")
                    
        parts.append("\n[RINGKASAN ELEMEN 2D]:")
        for k, v in sorted(entity_counts.items(), key=lambda x: -x[1]):
            parts.append(f"- {k}: {v} entitas")
            
        if texts:
            parts.append(f"\n[ANOTASI & ETIKET ({len(texts)})]:")
            for txt in texts[:20]:
                parts.append(f"- {txt}")
                
        return "\n".join(parts)[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing DXF {file_path}: {e}")
        return f"[Gagal membaca file CAD DXF: {e}]"

# ---------------------------------------------------------------------------
# 5. 3D Parametric CAD & SolidWorks (.step, .stp, .stl, .obj)
# ---------------------------------------------------------------------------
def parse_step_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts product definition, assembly hierarchy, and material metadata from ISO-10303 STEP (.step/.stp) files."""
    try:
        parts = ["=== ANALISIS MODEL CAD 3D / SOLIDWORKS (STEP/STP) ==="]
        products = set()
        materials = set()
        schemas = []
        author = "Tidak tertera"
        timestamp = "Tidak tertera"
        
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            for _ in range(5000):  # Read first 5000 lines for header & product definition
                line = f.readline()
                if not line:
                    break
                line_str = line.strip()
                
                # Header Author / Timestamp
                if "FILE_NAME" in line_str:
                    time_match = re.search(r"\'([^\']+)\'", line_str)
                    if time_match:
                        timestamp = time_match.group(1)
                if "FILE_SCHEMA" in line_str:
                    schema_match = re.findall(r"\'([^\']+)\'", line_str)
                    schemas.extend(schema_match)
                # Product & Assembly Name
                if "PRODUCT(" in line_str:
                    names = re.findall(r"\'([^\']+)\'", line_str)
                    for n in names:
                        if len(n) > 1 and not n.startswith("APPLICATION"):
                            products.add(n)
                # Material Designation
                if "MATERIAL_DESIGNATION" in line_str:
                    mats = re.findall(r"\'([^\']+)\'", line_str)
                    materials.update(mats)
                    
        parts.append(f"- Schema Standar: {', '.join(schemas) if schemas else 'ISO 10303-21 (AP203/AP214)'}")
        parts.append(f"- Waktu Pembuatan: {timestamp}")
        
        if products:
            parts.append(f"\n[STRUKTUR PART & ASSEMBLY ({len(products)})]:")
            for p in sorted(list(products))[:25]:
                parts.append(f"- Komponen: `{p}`")
        else:
            parts.append("- Model Part 3D Monolitik")
            
        if materials:
            parts.append(f"\n[SPESIFIKASI MATERIAL]:")
            for m in materials:
                parts.append(f"- Material: {m}")
                
        return "\n".join(parts)[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing STEP file {file_path}: {e}")
        return f"[Gagal membaca file CAD 3D STEP: {e}]"

def parse_stl_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts bounding box, volume, and facet mesh details from 3D STL / OBJ files."""
    try:
        import trimesh
        mesh = trimesh.load(file_path)
        parts = ["=== ANALISIS MODEL MESH 3D & 3D PRINTING (STL/OBJ) ==="]
        
        if isinstance(mesh, trimesh.Trimesh):
            bounds = mesh.bounding_box.extents
            parts.append(f"- Dimensi Bounding Box (P x L x T): {bounds[0]:.2f} x {bounds[1]:.2f} x {bounds[2]:.2f} mm")
            parts.append(f"- Luas Permukaan (Surface Area): {mesh.area:.2f} mm²")
            parts.append(f"- Estimasi Volume Padat: {mesh.volume:.2f} mm³ ({mesh.volume / 1000:.2f} cm³)")
            parts.append(f"- Jumlah Segitiga Mesh (Facets): {len(mesh.faces):,} triangles")
            parts.append(f"- Status Model Kedap Air (Watertight / 3D-Print Ready): {'Ya (Siap Cetak 3D)' if mesh.is_watertight else 'Tidak (Perlu Mesh Repair)'}")
        elif isinstance(mesh, trimesh.Scene):
            parts.append(f"- Model Assembly Multi-Mesh (Scene)")
            parts.append(f"- Total Objek Geometri: {len(mesh.geometry)} bagian")
            
        return "\n".join(parts)[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing STL mesh {file_path}: {e}")
        return f"[Gagal mengekstrak geometri 3D STL: {e}]"

# ---------------------------------------------------------------------------
# 6. CNC Machining & G-Code Files (.gcode, .nc, .cnc, .tap)
# ---------------------------------------------------------------------------
def parse_gcode_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts CNC machining commands, spindle speed, feed rate, and tools from G-Code / NC files."""
    try:
        parts = ["=== ANALISIS KODE PEMESINAN CNC (G-CODE / NC) ==="]
        total_lines = 0
        g_moves = 0
        spindles = set()
        feeds = set()
        tools = set()
        comments = []
        
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            for line in f:
                total_lines += 1
                line_str = line.strip()
                if not line_str:
                    continue
                    
                # Comments / Program Title
                if line_str.startswith("(") or line_str.startswith(";"):
                    if len(comments) < 10:
                        clean_c = line_str.strip("(); ")
                        if clean_c:
                            comments.append(clean_c)
                            
                # G0 / G1 / G2 / G3 moves
                if re.search(r'\b(G0|G1|G2|G3|G00|G01|G02|G03)\b', line_str, re.I):
                    g_moves += 1
                # Spindle Speed S
                s_match = re.search(r'\bS(\d+)\b', line_str)
                if s_match:
                    spindles.add(s_match.group(1))
                # Feed rate F
                f_match = re.search(r'\bF([\d\.]+)\b', line_str)
                if f_match:
                    feeds.add(f_match.group(1))
                # Tool T
                t_match = re.search(r'\bT(\d+)\b', line_str)
                if t_match:
                    tools.add(t_match.group(1))
                    
        parts.append(f"- Total Baris Perintah: {total_lines:,} baris")
        parts.append(f"- Total Langkah Pemotongan/Pergerakan Pahat: {g_moves:,} move commands")
        if tools:
            parts.append(f"- Nomor Pahat (Tool IDs): T{', T'.join(sorted(list(tools)))}")
        if spindles:
            parts.append(f"- Spindle Speed (RPM): {', '.join(sorted(list(spindles))[:5])} RPM")
        if feeds:
            parts.append(f"- Feed Rate (Kecepatan Pemakanan): {', '.join(sorted(list(feeds))[:5])} mm/min")
            
        if comments:
            parts.append("\n[KETERANGAN / HEADER PROGRAM CNC]:")
            for c in comments[:8]:
                parts.append(f"- {c}")
                
        return "\n".join(parts)[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing G-Code {file_path}: {e}")
        return f"[Gagal membaca file CNC G-Code: {e}]"

# ---------------------------------------------------------------------------
# 7. Discrete-Event Simulation Models (.fsm, .fsx - FlexSim)
# ---------------------------------------------------------------------------
def parse_flexsim_file(file_path: str, max_chars: int = MAX_TEXT_CHAR_LIMIT) -> str:
    """Extracts simulation entities, probability distributions, and logic from FlexSim .fsm and .fsx models."""
    try:
        ext = file_path.split(".")[-1].lower() if "." in file_path else ""
        
        # Handle XML model (.fsx)
        if ext == "fsx":
            tree = ET.parse(file_path)
            root = tree.getroot()
            parts = ["=== ANALISIS MODEL SIMULASI FLEXSIM (XML - .fsx) ==="]
            
            nodes = [elem.attrib.get("name", elem.tag) for elem in root.iter() if elem.attrib.get("name")]
            sources = [n for n in nodes if "source" in n.lower()]
            queues = [n for n in nodes if "queue" in n.lower() or "buffer" in n.lower()]
            processors = [n for n in nodes if "processor" in n.lower() or "machine" in n.lower() or "station" in n.lower()]
            sinks = [n for n in nodes if "sink" in n.lower()]
            conveyors = [n for n in nodes if "conveyor" in n.lower()]
            
            parts.append(f"- Pembangkit Entitas (Source): {', '.join(set(sources)) if sources else '1 Unit'}")
            parts.append(f"- Buffer / Antrean (Queue): {', '.join(set(queues)) if queues else 'Tidak ada'}")
            parts.append(f"- Mesin Produksi (Processor): {', '.join(set(processors)) if processors else 'Tidak ada'}")
            parts.append(f"- Alur Keluar (Sink): {', '.join(set(sinks)) if sinks else '1 Unit'}")
            parts.append(f"- Conveyor / Jalur: {', '.join(set(conveyors)) if conveyors else 'Direct Transfer'}")
            return "\n".join(parts)[:max_chars]
            
        # Handle Binary model (.fsm)
        with open(file_path, "rb") as f:
            raw = f.read()

        gzip_offset = raw.find(b"\x1f\x8b\x08")
        if gzip_offset == -1:
            data = raw
        else:
            data = gzip.decompress(raw[gzip_offset:])

        # Extract strings
        strings = re.findall(rb'[A-Za-z0-9_ \-\.\(\)\,\:\/\\><\=\+\*\[\]\"\'\$\@\#]{3,}', data)
        decoded = [s.decode('latin1', errors='ignore').strip() for s in strings if len(s.strip()) >= 3]

        # Extract components
        sources = sorted(list(set(s for s in decoded if re.match(r'^(Source|Zrodlo|Generator)\w*$', s, re.I))))
        queues = sorted(list(set(s for s in decoded if re.match(r'^(Queue|Kolejka|Bufor|Buffer)\w*$', s, re.I))))
        processors = sorted(list(set(s for s in decoded if re.match(r'^(Processor|Maszyna|Stacja|Station|CNC|Milling|Lathe|Drill)\w*$', s, re.I))))
        sinks = sorted(list(set(s for s in decoded if re.match(r'^(Sink|Koniec|Ujscie|Exit)\w*$', s, re.I))))
        conveyors = sorted(list(set(s for s in decoded if re.match(r'^(Conveyor|Przenosnik|Podajnik)\w*$', s, re.I))))
        operators = sorted(list(set(s for s in decoded if re.match(r'^(Operator|Robot|Pracownik|Worker|AGV|Transporter)\w*$', s, re.I))))
        racks = sorted(list(set(s for s in decoded if re.match(r'^(Rack|Magazyn|Storage|ASRS|Warehouse)\w*$', s, re.I))))

        # Statistical distributions
        dists = list(set(s for s in decoded if any(d in s.lower() for d in ['exponential(', 'triangular(', 'normal(', 'uniform(', 'duniform(', 'lognormal(', 'weibull(', 'bernoulli(', 'poisson('])))

        # Logic / Triggers
        triggers = list(set(s for s in decoded if any(k in s.lower() for k in ['eventlisten', 'onentry', 'onexit', 'onreset', 'onsetup', 'processflow', 'transportdispatcher'])))

        summary = [
            "=== ANALISIS MODEL SIMULASI FLEXSIM (.fsm) ===",
            f"- Pembangkit Entitas (Source): {', '.join(sources) if sources else '1 Unit'}",
            f"- Buffer / Antrean (Queue): {', '.join(queues) if queues else 'Tidak ada'}",
            f"- Stasiun Mesin (Processor): {', '.join(processors) if processors else 'Tidak ada'}",
            f"- Alur Keluar (Sink): {', '.join(sinks) if sinks else '1 Unit'}",
            f"- Material Handling (Conveyor): {', '.join(conveyors) if conveyors else 'Direct Transfer'}",
            f"- Operator & AGV: {', '.join(operators) if operators else 'Automated'}",
            f"- Penyimpanan & Gudang (Rack/ASRS): {', '.join(racks) if racks else 'Tidak ada'}",
            "",
            "=== DISTRIBUSI WAKTU & PARAMETER PROSES ==="
        ]
        if dists:
            for d in dists[:10]:
                summary.append(f"- {d}")
        else:
            summary.append("- Parameter waktu proses terdefinisi pada node tabel.")

        if triggers:
            summary.append("\n=== LOGIKA PROSES & EVENT TRIGGERS ===")
            for t in triggers[:8]:
                summary.append(f"- `{t}`")

        return "\n".join(summary)[:max_chars]
    except Exception as e:
        logger.error(f"Error parsing FlexSim model {file_path}: {e}")
        return f"[Gagal membaca model FlexSim: {e}]"

# ---------------------------------------------------------------------------
# Universal Dispatcher
# ---------------------------------------------------------------------------
def extract_document_content(file_path: str, original_filename: str) -> str:
    """Universal dispatcher to parse any supported document, CAD, simulation, or code file."""
    if not os.path.exists(file_path):
        return f"[Berkas {original_filename} tidak ditemukan di server]"
        
    ext = original_filename.split(".")[-1].lower() if "." in original_filename else ""
    
    # 1. Office & Text Docs
    if ext in ["docx"]:
        parsed = parse_docx_file(file_path)
    elif ext in ["pptx", "ppt"]:
        parsed = parse_pptx_file(file_path)
    elif ext in ["xlsx", "xls"]:
        parsed = parse_excel_file(file_path)
    elif ext in ["csv"]:
        parsed = parse_csv_file(file_path)
    elif ext in ["pdf"]:
        parsed = parse_pdf_file(file_path)
    elif ext in ["zip", "tar", "gz", "7z"]:
        parsed = parse_zip_file(file_path)
        
    # 2. AutoCAD & 2D CAD
    elif ext in ["dwg"]:
        parsed = parse_dwg_file(file_path)
    elif ext in ["dxf"]:
        parsed = parse_dxf_file(file_path)
        
    # 3. 3D CAD & SolidWorks
    elif ext in ["step", "stp"]:
        parsed = parse_step_file(file_path)
    elif ext in ["stl", "obj"]:
        parsed = parse_stl_file(file_path)
        
    # 4. CNC Machining
    elif ext in ["gcode", "nc", "cnc", "tap"]:
        parsed = parse_gcode_file(file_path)
        
    # 5. Discrete Event Simulation (FlexSim)
    elif ext in ["fsm", "fsx"]:
        parsed = parse_flexsim_file(file_path)
        
    # 6. Source Code & Scripts
    elif ext in [
        "py", "js", "ts", "tsx", "jsx", "html", "css", "json", "yaml", "yml",
        "sql", "sh", "bash", "c", "cpp", "h", "hpp", "java", "kt", "rs",
        "go", "php", "rb", "r", "m", "txt", "md", "markdown", "env", "log"
    ]:
        content = parse_text_file(file_path)
        parsed = f"```{ext}\n{content}\n```"
    else:
        content = parse_text_file(file_path)
        parsed = content if content else f"[Format berkas .{ext} disimpan sebagai lampiran biner]"
        
    return f"\n\n=== LAMPIRAN DOKUMEN/BERKAS: {original_filename} ===\n{parsed}\n"
