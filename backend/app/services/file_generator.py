import os
import re
import uuid
import logging
from typing import List, Dict, Any, Optional

logger = logging.getLogger("RuangTI.FileGenerator")

GENERATED_FILES_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "uploads", "generated")
os.makedirs(GENERATED_FILES_DIR, exist_ok=True)

# ---------------------------------------------------------------------------
# 1. Excel Spreadsheet Generator (.xlsx)
# ---------------------------------------------------------------------------
def generate_excel_file(
    filename: str,
    title: str = "Tabel Data RuangTI",
    headers: Optional[List[str]] = None,
    rows: Optional[List[List[Any]]] = None,
    sheet_name: str = "Sheet1"
) -> str:
    """Generates a professionally styled Excel workbook."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = str(sheet_name)[:31]
    ws.views.sheetView[0].showGridLines = True

    # Color Tokens
    HEADER_FILL = PatternFill(start_color="1C1F26", end_color="1C1F26", fill_type="solid")
    ACCENT_FILL = PatternFill(start_color="E09F3E", end_color="E09F3E", fill_type="solid")
    ZEBRA_FILL = PatternFill(start_color="F8F9FA", end_color="F8F9FA", fill_type="solid")
    HEADER_FONT = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
    TITLE_FONT = Font(name="Segoe UI", size=14, bold=True, color="16181D")
    DATA_FONT = Font(name="Segoe UI", size=10, color="16181D")
    BOLD_FONT = Font(name="Segoe UI", size=10, bold=True, color="16181D")

    # Borders
    THIN_GRAY = Side(border_style="thin", color="E2E4E9")
    CELL_BORDER = Border(left=THIN_GRAY, right=THIN_GRAY, top=THIN_GRAY, bottom=THIN_GRAY)
    DOUBLE_BOTTOM = Border(left=THIN_GRAY, right=THIN_GRAY, top=THIN_GRAY, bottom=Side(border_style="double", color="16181D"))

    # Title Banner
    current_row = 1
    if title:
        ws.cell(row=current_row, column=1, value=title).font = TITLE_FONT
        current_row += 2

    # Headers
    if headers:
        for col_idx, h_text in enumerate(headers, 1):
            cell = ws.cell(row=current_row, column=col_idx, value=str(h_text).strip())
            cell.fill = HEADER_FILL
            cell.font = HEADER_FONT
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = CELL_BORDER
        ws.row_dimensions[current_row].height = 28
        current_row += 1

    # Data Rows
    if rows:
        for r_idx, row_data in enumerate(rows):
            is_zebra = r_idx % 2 == 1
            ws.row_dimensions[current_row].height = 20
            
            for c_idx, val in enumerate(row_data, 1):
                cell = ws.cell(row=current_row, column=c_idx)
                
                # Auto-parse numeric / float / currency
                val_str = str(val).strip() if val is not None else ""
                clean_num = val_str.replace(",", ".").replace("Rp", "").replace("%", "").strip()
                
                if val_str.endswith("%") and clean_num.replace(".", "", 1).isdigit():
                    cell.value = float(clean_num) / 100.0
                    cell.number_format = "0.0%"
                    cell.alignment = Alignment(horizontal="right", vertical="center")
                elif clean_num.replace(".", "", 1).replace("-", "", 1).isdigit() and len(clean_num) > 0 and not clean_num.startswith("0") or clean_num == "0":
                    try:
                        num_val = float(clean_num) if "." in clean_num else int(clean_num)
                        cell.value = num_val
                        cell.number_format = "#,##0" if isinstance(num_val, int) else "#,##0.00"
                        cell.alignment = Alignment(horizontal="right", vertical="center")
                    except ValueError:
                        cell.value = val_str
                        cell.alignment = Alignment(horizontal="left", vertical="center")
                else:
                    cell.value = val_str
                    cell.alignment = Alignment(horizontal="left", vertical="center")

                cell.font = DATA_FONT
                cell.border = CELL_BORDER
                if is_zebra:
                    cell.fill = ZEBRA_FILL
            current_row += 1

    # Auto-fit Column Widths (Anti Clipping)
    for col in ws.columns:
        max_len = 0
        col_letter = get_column_letter(col[0].column)
        for cell in col:
            if cell.row == 1 and title:
                continue
            val_s = str(cell.value or "")
            if len(val_s) > max_len:
                max_len = len(val_s)
        ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

    # Save to generated directory
    safe_name = re.sub(r'[^a-zA-Z0-9_\-\.]', '_', filename)
    if not safe_name.endswith(".xlsx"):
        safe_name += ".xlsx"
        
    out_path = os.path.join(GENERATED_FILES_DIR, f"{uuid.uuid4().hex[:8]}_{safe_name}")
    wb.save(out_path)
    logger.info(f"Generated professional Excel file: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# 2. Word Document Generator (.docx)
# ---------------------------------------------------------------------------
def generate_docx_file(
    filename: str,
    title: str = "Dokumen RuangTI",
    sections: Optional[List[Dict[str, Any]]] = None
) -> str:
    """Generates a clean, professionally formatted Microsoft Word (.docx) document."""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    doc = Document()
    
    # Page Margins (Normal 1 inch)
    for s in doc.sections:
        s.top_margin = Inches(1)
        s.bottom_margin = Inches(1)
        s.left_margin = Inches(1)
        s.right_margin = Inches(1)

    # Document Header Title
    if title:
        title_p = doc.add_paragraph()
        title_p.paragraph_format.space_before = Pt(0)
        title_p.paragraph_format.space_after = Pt(14)
        run = title_p.add_run(title)
        run.font.name = "Segoe UI"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(22, 24, 29)

    if sections:
        for sec in sections:
            heading = sec.get("heading")
            if heading:
                h_p = doc.add_paragraph()
                h_p.paragraph_format.space_before = Pt(14)
                h_p.paragraph_format.space_after = Pt(4)
                h_run = h_p.add_run(heading)
                h_run.font.name = "Segoe UI"
                h_run.font.size = Pt(13)
                h_run.font.bold = True
                h_run.font.color.rgb = RGBColor(224, 159, 62)  # Amber Gold Heading

            # Paragraphs
            paragraphs = sec.get("paragraphs", [])
            for p_text in paragraphs:
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(6)
                p.paragraph_format.line_spacing = 1.15
                run = p.add_run(str(p_text))
                run.font.name = "Segoe UI"
                run.font.size = Pt(10.5)
                run.font.color.rgb = RGBColor(50, 55, 65)

            # Bullet points
            bullets = sec.get("bullets", [])
            for b_text in bullets:
                bp = doc.add_paragraph(style="List Bullet")
                bp.paragraph_format.space_after = Pt(3)
                run = bp.add_run(str(b_text))
                run.font.name = "Segoe UI"
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(50, 55, 65)

            # Table in Section
            table_data = sec.get("table")
            if table_data and isinstance(table_data, dict):
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers:
                    t = doc.add_table(rows=len(rows) + 1, cols=len(headers))
                    t.alignment = WD_TABLE_ALIGNMENT.CENTER
                    
                    # Style Header
                    for c_idx, h in enumerate(headers):
                        cell = t.cell(0, c_idx)
                        cell.text = str(h)
                        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="1C1F26"/>')
                        cell._tc.get_or_add_tcPr().append(shading)
                        for cp in cell.paragraphs:
                            for cr in cp.runs:
                                cr.font.name = "Segoe UI"
                                cr.font.size = Pt(9.5)
                                cr.font.bold = True
                                cr.font.color.rgb = RGBColor(255, 255, 255)
                                
                    # Style Data Rows
                    for r_idx, row_vals in enumerate(rows, 1):
                        bg_color = "F8F9FA" if r_idx % 2 == 0 else "FFFFFF"
                        for c_idx, val in enumerate(row_vals):
                            if c_idx < len(headers):
                                cell = t.cell(r_idx, c_idx)
                                cell.text = str(val if val is not None else "")
                                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{bg_color}"/>')
                                cell._tc.get_or_add_tcPr().append(shading)
                                for cp in cell.paragraphs:
                                    for cr in cp.runs:
                                        cr.font.name = "Segoe UI"
                                        cr.font.size = Pt(9)
                                        cr.font.color.rgb = RGBColor(30, 35, 45)

    safe_name = re.sub(r'[^a-zA-Z0-9_\-\.]', '_', filename)
    if not safe_name.endswith(".docx"):
        safe_name += ".docx"
        
    out_path = os.path.join(GENERATED_FILES_DIR, f"{uuid.uuid4().hex[:8]}_{safe_name}")
    doc.save(out_path)
    logger.info(f"Generated professional Word document: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# 3. PowerPoint Presentation Generator (.pptx)
# ---------------------------------------------------------------------------
def generate_pptx_file(
    filename: str,
    title: str = "Presentasi RuangTI",
    subtitle: str = "Industrial Engineering Workspace",
    slides_data: Optional[List[Dict[str, Any]]] = None
) -> str:
    """Generates a modern 16:9 widescreen PowerPoint presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    DARK_BG = RGBColor(28, 31, 38)
    GOLD = RGBColor(224, 159, 62)
    WHITE = RGBColor(255, 255, 255)
    GRAY_TEXT = RGBColor(160, 165, 175)
    CARD_BG = RGBColor(248, 249, 250)
    CARD_TEXT = RGBColor(30, 35, 45)

    # 1. Slide 1: Cover Slide
    cover_slide = prs.slides.add_slide(blank_layout)
    bg_shape = cover_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BG
    bg_shape.line.fill.background()

    # Accent Top Bar
    bar = cover_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(0.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    # Cover Title Box
    title_box = cover_slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(18)
        p2.font.color.rgb = GOLD
        p2.space_before = Pt(12)

    # 2. Content Slides
    if slides_data:
        for s_idx, s_info in enumerate(slides_data, 1):
            slide = prs.slides.add_slide(blank_layout)
            
            # Slide Header Bar
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
            htf = header_box.text_frame
            hp = htf.paragraphs[0]
            hp.text = s_info.get("title", f"Topik Pembahasan {s_idx}")
            hp.font.name = "Segoe UI"
            hp.font.size = Pt(24)
            hp.font.bold = True
            hp.font.color.rgb = DARK_BG

            # Slide Content / Cards
            points = s_info.get("points", [])
            if points:
                box_y = Inches(1.8)
                box_w = Inches(11.7)
                box_h = Inches(4.8)
                
                content_box = slide.shapes.add_textbox(Inches(0.8), box_y, box_w, box_h)
                ctf = content_box.text_frame
                ctf.word_wrap = True
                
                for p_idx, pt_text in enumerate(points):
                    cp = ctf.paragraphs[0] if p_idx == 0 else ctf.add_paragraph()
                    cp.text = f"•  {pt_text}"
                    cp.font.name = "Segoe UI"
                    cp.font.size = Pt(16)
                    cp.font.color.rgb = CARD_TEXT
                    cp.space_before = Pt(14)
                    cp.line_spacing = 1.25

    safe_name = re.sub(r'[^a-zA-Z0-9_\-\.]', '_', filename)
    if not safe_name.endswith(".pptx"):
        safe_name += ".pptx"
        
    out_path = os.path.join(GENERATED_FILES_DIR, f"{uuid.uuid4().hex[:8]}_{safe_name}")
    prs.save(out_path)
    logger.info(f"Generated professional PowerPoint presentation: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# 4. PDF Document Generator (.pdf)
# ---------------------------------------------------------------------------
def generate_pdf_file(
    filename: str,
    title: str = "Laporan Resmi RuangTI",
    sections: Optional[List[Dict[str, Any]]] = None
) -> str:
    """Generates a publication-grade PDF document using ReportLab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    safe_name = re.sub(r'[^a-zA-Z0-9_\-\.]', '_', filename)
    if not safe_name.endswith(".pdf"):
        safe_name += ".pdf"
        
    out_path = os.path.join(GENERATED_FILES_DIR, f"{uuid.uuid4().hex[:8]}_{safe_name}")

    doc = SimpleDocTemplate(
        out_path,
        pagesize=A4,
        rightMargin=40,
        leftMargin=40,
        topMargin=40,
        bottomMargin=40
    )

    styles = getSampleStyleSheet()
    
    # Custom Palette
    COLOR_PRIMARY = colors.HexColor("#16181D")
    COLOR_ACCENT = colors.HexColor("#E09F3E")
    COLOR_BG_HEADER = colors.HexColor("#1C1F26")
    COLOR_MUTED = colors.HexColor("#5F6570")

    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=COLOR_PRIMARY,
        spaceAfter=12
    )

    h2_style = ParagraphStyle(
        'DocHeading',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=16,
        textColor=COLOR_ACCENT,
        spaceBefore=14,
        spaceAfter=6
    )

    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        textColor=COLOR_PRIMARY,
        spaceAfter=6
    )

    story = []
    
    # Title
    if title:
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 8))

    if sections:
        for sec in sections:
            heading = sec.get("heading")
            if heading:
                story.append(Paragraph(heading, h2_style))
                
            for p_text in sec.get("paragraphs", []):
                story.append(Paragraph(str(p_text), body_style))

            bullets = sec.get("bullets", [])
            for b_text in bullets:
                story.append(Paragraph(f"• {b_text}", body_style))

            table_data = sec.get("table")
            if table_data and isinstance(table_data, dict):
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers:
                    t_rows = [[Paragraph(f"<b>{h}</b>", ParagraphStyle('TH', fontName='Helvetica-Bold', fontSize=9, textColor=colors.white)) for h in headers]]
                    for row in rows:
                        t_rows.append([Paragraph(str(v if v is not None else ""), ParagraphStyle('TD', fontName='Helvetica', fontSize=8.5, leading=11)) for v in row])
                        
                    pdf_table = Table(t_rows, repeatRows=1)
                    pdf_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), COLOR_BG_HEADER),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('INNERGRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#E2E4E9")),
                        ('BOX', (0, 0), (-1, -1), 0.5, colors.HexColor("#D0D3D9")),
                        ('TOPPADDING', (0, 0), (-1, -1), 4),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                    ]))
                    story.append(Spacer(1, 4))
                    story.append(pdf_table)
                    story.append(Spacer(1, 8))

    doc.build(story)
    logger.info(f"Generated professional PDF document: {out_path}")
    return out_path
